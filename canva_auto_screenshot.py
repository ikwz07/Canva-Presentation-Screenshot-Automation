from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.chrome.options import Options
import pyautogui
from pptx import Presentation
from pptx.util import Inches
import time
import os

# ==== CONFIGURATION ====
CANVA_LINK = 'https://www.canva.com/design/DAGlEitHo3E/5-AvJk899DX1bNFsJ34doA/view?utm_content=DAGlEitHo3E&utm_campaign=designshare&utm_medium=link2&utm_source=uniquelinks&utlId=h50ac65cb4c'
TOTAL_SLIDES = 10  # Set the total number of slides you want to capture
SCREENSHOT_DIR = 'screenshots'
PPTX_FILENAME = 'Canva_Presentation_Auto.pptx'
CHROMEDRIVER_PATH = 'C:\\Users\\Imran Khan WZ\\PycharmProjects\\CanvaAutoScreenShotsProject\\chromedriver-win64\\chromedriver.exe'# Update if needed
SCREEN_WIDTH = 2560
SCREEN_HEIGHT = 1440

# Create directory for screenshots
if not os.path.exists(SCREENSHOT_DIR):
    os.makedirs(SCREENSHOT_DIR)

# Set up Selenium with Chrome
chrome_options = Options()
chrome_options.add_argument("--start-maximized")
service = Service(CHROMEDRIVER_PATH)
driver = webdriver.Chrome(service=service, options=chrome_options)

# Open Canva link
driver.get(CANVA_LINK)
time.sleep(5)  # Wait for the page to load completely

# Click the 'Present' button
pyautogui.press('f')


# Wait for presentation to load
print("Waiting for presentation mode to start...")
time.sleep(7)

# Take screenshots and navigate slides
for i in range(TOTAL_SLIDES):
    filename = f"{SCREENSHOT_DIR}/slide_{i + 1}.png"
    screenshot = pyautogui.screenshot(region=(0, 0, SCREEN_WIDTH, SCREEN_HEIGHT))
    screenshot.save(filename)
    print(f"Captured {filename}")

    # Press right arrow key to go to next slide (skip on last slide)
    if i < TOTAL_SLIDES - 1:
        pyautogui.press('right')
        time.sleep(3)

# Close the browser
driver.quit()

# Create PowerPoint presentation
prs = Presentation()
slide_width = prs.slide_width
slide_height = prs.slide_height

for i in range(TOTAL_SLIDES):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
    img_path = f"{SCREENSHOT_DIR}/slide_{i + 1}.png"
    slide.shapes.add_picture(img_path, 0, 0, width=slide_width, height=slide_height)

# Save the presentation
prs.save(PPTX_FILENAME)
print(f"Presentation created: {PPTX_FILENAME}") 

